import io
import base64
import logging
from typing import List
from pptx import Presentation
from PIL import Image

from .base import ExtractedPage
from providers.client import NeuroFlowClient
from providers.router import RoutingCriteria
from providers.base import ChatMessage

logger = logging.getLogger(__name__)

class PPTXExtractor:
    """Extractor for PPTX files.
    
    Extracts text from shapes and speaker notes.
    If a slide contains an image, it extracts the image and queries a vision LLM
    for a description of the diagram.
    """
    
    async def extract(self, file_path_or_bytes: str | bytes, **kwargs) -> list[ExtractedPage]:
        pages: list[ExtractedPage] = []
        
        if isinstance(file_path_or_bytes, bytes):
            ppt_file = io.BytesIO(file_path_or_bytes)
        else:
            ppt_file = file_path_or_bytes
            
        try:
            prs = Presentation(ppt_file)
        except Exception as e:
            logger.error(f"Failed to load PPTX: {e}")
            return pages

        client = NeuroFlowClient()
            
        for i, slide in enumerate(prs.slides):
            page_num = i + 1
            
            slide_text = []
            images = []
            
            # Extract text and images from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
                
                # Check for images
                if hasattr(shape, "image"):
                    try:
                        img_bytes = shape.image.blob
                        img = Image.open(io.BytesIO(img_bytes))
                        if img.mode not in ("RGB", "L"):
                            img = img.convert("RGB")
                        images.append(img)
                    except Exception as e:
                        logger.warning(f"Failed to extract image from slide {page_num}: {e}")

            # Extract speaker notes
            if slide.has_notes_slide:
                text_frame = slide.notes_slide.notes_text_frame
                if text_frame and text_frame.text.strip():
                    slide_text.append(f"Speaker Notes:\n{text_frame.text.strip()}")
            
            # Combine text
            combined_text = "\n".join(slide_text)
            
            # Process images with Vision LLM
            vision_descriptions = []
            for img_idx, img in enumerate(images):
                try:
                    # Resize if too large
                    max_size = 1024
                    if img.width > max_size or img.height > max_size:
                        ratio = min(max_size / img.width, max_size / img.height)
                        new_size = (int(img.width * ratio), int(img.height * ratio))
                        img = img.resize(new_size, Image.Resampling.LANCZOS)
                        
                    buffered = io.BytesIO()
                    img.save(buffered, format="JPEG")
                    b64_img = base64.b64encode(buffered.getvalue()).decode("utf-8")
                    
                    content = [
                        {"type": "text", "text": "Describe this slide diagram/image in detail."},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64_img}"}}
                    ]
                    
                    messages = [ChatMessage(role="user", content=content)]
                    criteria = RoutingCriteria(require_vision=True, task_type="rag_generation")
                    
                    result = await client.chat(messages, routing_criteria=criteria)
                    vision_descriptions.append(f"Image {img_idx+1} Description: {result.content.strip()}")
                except Exception as e:
                    logger.error(f"Vision LLM failed for slide {page_num} image: {e}")
            
            if vision_descriptions:
                combined_text += "\n\n" + "\n\n".join(vision_descriptions)
                
            pages.append(
                ExtractedPage(
                    page_number=page_num,
                    content=combined_text,
                    content_type="text",
                    metadata={"has_images": len(images) > 0}
                )
            )
            
        return pages
